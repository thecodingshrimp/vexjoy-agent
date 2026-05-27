#!/usr/bin/env python3
"""Unified deck pipeline entry point.

Single-command CLI that runs:
    HTML deck  ->  slide-map JSON  ->  .pptx (and optional QA PNGs)

Usage:
    python3 run-unified.py --input deck.html --format pptx --out ./out/
    python3 run-unified.py --input deck.html --format pptx --out deck.pptx

`--out` accepts either a directory (sibling artifacts written next to the
.pptx) or a `.pptx` file path (single-file mode; no sibling JSON/report).

Format support:
    pptx  -> implemented (this script)
    pdf   -> NOT WIRED here; use scripts/to-pdf.py for HTML->PDF.

Exit codes:
    0  ok
    2  bad input / missing tool
    3  conversion failure
"""

from __future__ import annotations

import argparse
import json
import shutil
import subprocess
import sys
from pathlib import Path

HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(HERE))

# Local imports from this directory
import _pptx_engine as pptx_engine
import extract_slides

# ---------------------------------------------------------------------------
# Steps
# ---------------------------------------------------------------------------


def step_extract(html_path: Path, json_out: Path) -> list[dict]:
    """Step 1: HTML -> slide-map JSON."""
    html = html_path.read_text(encoding="utf-8")
    slides = extract_slides.extract_slides(html)
    json_out.parent.mkdir(parents=True, exist_ok=True)
    json_out.write_text(json.dumps(slides, indent=2), encoding="utf-8")
    return slides


def step_build_pptx(slides: list[dict], pptx_out: Path) -> Path:
    """Step 2: slide-map JSON -> .pptx via copied pptx engine."""
    design = {"palette": "minimal"}
    pptx_out.parent.mkdir(parents=True, exist_ok=True)
    pptx_engine.build_presentation(slides, design, str(pptx_out))
    return pptx_out


def step_render_qa(pptx_path: Path, render_dir: Path) -> tuple[bool, str]:
    """Step 3 (optional): .pptx -> per-slide PNGs via LibreOffice.

    Returns (rendered, message). Skips silently if soffice is unavailable.
    """
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return False, "skipped (soffice not on PATH)"
    render_dir.mkdir(parents=True, exist_ok=True)
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(render_dir), str(pptx_path)],
            check=True,
            capture_output=True,
            timeout=120,
        )
        pdfs = list(render_dir.glob("*.pdf"))
        if not pdfs:
            return False, "soffice produced no PDF"
        return True, f"PDF at {pdfs[0]}"
    except Exception as exc:  # pragma: no cover - host-dependent
        return False, f"soffice failed: {exc}"


# ---------------------------------------------------------------------------
# Fidelity report
# ---------------------------------------------------------------------------


def step_report(
    pptx_path: Path,
    slides: list[dict],
    report_path: Path,
    rendered_msg: str,
) -> dict:
    """Open the generated .pptx with python-pptx, compare to baseline, write report.md."""
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    slide_count = len(prs.slides)
    text_frames = sum(1 for sl in prs.slides for sh in sl.shapes if sh.has_text_frame)
    width_in = prs.slide_width / 914400.0
    height_in = prs.slide_height / 914400.0
    aspect = width_in / height_in
    is_widescreen = abs(aspect - 16 / 9) < 0.02
    file_size = pptx_path.stat().st_size

    # Type distribution
    type_counts: dict[str, int] = {}
    for s in slides:
        t = s.get("type", "?")
        type_counts[t] = type_counts.get(t, 0) + 1

    # Baseline comparison (optional)
    baseline_dir = Path(".audit/pptx-test/render")
    baseline_pngs = sorted(baseline_dir.glob("original-*.png")) if baseline_dir.exists() else []

    # ---------- Score ----------
    # Axes (each 0-2):
    #   slide_count_match      (12 expected)
    #   text_frame_density     (>=120 strong, 60-119 partial, <60 weak)
    #   aspect_ratio_widescreen
    #   no_unrecognized_types  (all extracted types present)
    #   build_succeeded        (file exists, parsable)
    score = 0
    worst_axis = None
    worst_score = 99
    axes = []

    def add_axis(name: str, points: int, reason: str) -> None:
        nonlocal score, worst_axis, worst_score
        score += points
        axes.append((name, points, reason))
        if points < worst_score:
            worst_score = points
            worst_axis = name

    add_axis(
        "slide_count_match",
        2 if slide_count == 12 else (1 if abs(slide_count - 12) <= 2 else 0),
        f"got {slide_count}, expected 12",
    )
    add_axis(
        "text_frame_density",
        2 if text_frames >= 120 else (1 if text_frames >= 60 else 0),
        f"got {text_frames}, expected >=120",
    )
    add_axis(
        "aspect_ratio_widescreen",
        2 if is_widescreen else 0,
        f"{width_in:.2f}x{height_in:.2f}in (aspect {aspect:.3f})",
    )
    # Types in extractor that fall back to default in engine = layout fidelity loss.
    # Pulled from the engine's registry so the two stay in sync.
    supported = set(getattr(pptx_engine, "SUPPORTED_LAYOUTS", set()))
    if not supported:  # pragma: no cover - belt-and-suspenders
        supported = {"title", "content", "closing"}

    # Normalize extractor types the same way the engine does before lookup.
    def _norm(t: str) -> str:
        return str(t).lower().strip().replace("-", "_").replace(" ", "_")

    fallback_types = [t for t in type_counts if _norm(t) not in supported]
    fallback_slides = sum(type_counts[t] for t in fallback_types)
    add_axis(
        "layout_coverage",
        2 if not fallback_types else (1 if fallback_slides <= 6 else 0),
        f"{fallback_slides}/{slide_count} slides fall back to default layout: {fallback_types}",
    )
    add_axis(
        "build_succeeded",
        2 if file_size > 10_000 else 0,
        f"{file_size:,} bytes",
    )

    # ---------- Write report ----------
    lines = [
        "# Unified Deck Prototype — V1 Fidelity Report",
        "",
        f"- **Source HTML**: {pptx_path.name.replace('.pptx', '.html')}",
        f"- **Generated PPTX**: `{pptx_path}`",
        f"- **Size**: {file_size:,} bytes",
        f"- **QA render**: {rendered_msg}",
        "",
        "## Slide-level metrics",
        "",
        "| Metric | Value | Expected |",
        "|---|---|---|",
        f"| Slide count | {slide_count} | 12 |",
        f"| Text frames (total) | {text_frames} | >=120 |",
        f"| Slide dimensions | {width_in:.2f} x {height_in:.2f} in | 13.33 x 7.50 in (16:9) |",
        f"| Aspect ratio | {aspect:.3f} | 1.778 |",
        "",
        "## Type distribution (from extractor)",
        "",
        "| Type | Count | Engine handling |",
        "|---|---|---|",
    ]
    for t, c in sorted(type_counts.items(), key=lambda kv: -kv[1]):
        handling = "native" if _norm(t) in supported else "FALLBACK -> content"
        lines.append(f"| `{t}` | {c} | {handling} |")
    lines += [
        "",
        "## Fidelity score",
        "",
        "| Axis | Points (/2) | Detail |",
        "|---|---|---|",
    ]
    for n, p, r in axes:
        lines.append(f"| {n} | {p} | {r} |")
    total_max = 2 * len(axes)
    lines += [
        "",
        f"**Total: {score} / {total_max}** (worst axis: `{worst_axis}` at {worst_score}/2)",
        "",
        "## Baseline comparison",
        "",
    ]
    if baseline_pngs:
        lines.append(f"Baseline PNGs found: {len(baseline_pngs)} files in `{baseline_dir}/`.")
        lines.append("Per-slide perceptual diff is not implemented in V1; visual comparison")
        lines.append("requires manual side-by-side review or a future image-diff axis.")
    else:
        lines.append("No baseline PNGs found at `.audit/pptx-test/render/` — comparison skipped.")
    lines += [
        "",
        "## Known gaps (Phase 2 scope)",
        "",
        "- PDF output (`--format pdf`) NOT wired; would require LibreOffice headless or weasyprint.",
        "- Perceptual diff vs baseline PNGs not implemented (would require pixelmatch / SSIM).",
        "- LibreOffice render step skipped on hosts without `soffice`.",
        "- Font embedding not implemented; relies on Aptos / Cascadia Code being present",
        "  on the viewer machine (PowerPoint 2023+ ships them).",
    ]
    report_path.write_text("\n".join(lines) + "\n", encoding="utf-8")

    # Normalize fidelity to a 0-10 score
    fidelity_10 = round(score / total_max * 10, 1)
    return {
        "score": score,
        "score_max": total_max,
        "fidelity_10": fidelity_10,
        "worst_axis": worst_axis,
        "worst_score": worst_score,
        "slide_count": slide_count,
        "text_frames": text_frames,
        "fallback_types": fallback_types,
    }


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def main() -> int:
    parser = argparse.ArgumentParser(description="Unified deck pipeline (HTML -> PPTX)")
    parser.add_argument("--input", required=True, help="HTML deck file")
    parser.add_argument(
        "--format",
        choices=["pptx", "pdf"],
        default="pptx",
        help="Output format (pdf NOT wired here; use to-pdf.py)",
    )
    parser.add_argument("--out", required=True, help="Output directory or .pptx file path")
    parser.add_argument(
        "--no-render",
        action="store_true",
        help="Skip optional LibreOffice QA render step",
    )
    args = parser.parse_args()

    if args.format == "pdf":
        print("FAIL: --format pdf is not wired here; use scripts/to-pdf.py", file=sys.stderr)
        return 2

    html_path = Path(args.input).resolve()
    if not html_path.exists():
        print(f"FAIL: input HTML not found: {html_path}", file=sys.stderr)
        return 2

    out_arg = Path(args.out).resolve()
    # Single-file mode: --out points at a .pptx (parent dir is the working dir).
    # Directory mode: --out is (or will become) a directory; sibling artifacts
    # land next to the .pptx.
    if out_arg.suffix.lower() == ".pptx":
        out_dir = out_arg.parent
        pptx_path = out_arg
        single_file_mode = True
    else:
        out_dir = out_arg
        pptx_path = out_dir / (html_path.stem + ".pptx")
        single_file_mode = False
    out_dir.mkdir(parents=True, exist_ok=True)

    json_path = out_dir / "slides.json"
    report_path = out_dir / "report.md"
    render_dir = out_dir / "render"

    try:
        # Step 1
        slides = step_extract(html_path, json_path)
        if not slides:
            print("FAIL: extractor produced 0 slides", file=sys.stderr)
            return 3
        # Step 2
        step_build_pptx(slides, pptx_path)
        if not pptx_path.exists() or pptx_path.stat().st_size < 1000:
            print(f"FAIL: pptx not produced at {pptx_path}", file=sys.stderr)
            return 3
        # Step 3 (optional)
        if args.no_render:
            rendered_msg = "skipped (--no-render)"
        else:
            ok, msg = step_render_qa(pptx_path, render_dir)
            rendered_msg = msg if ok else f"skipped ({msg})"
        # Step 4
        result = step_report(pptx_path, slides, report_path, rendered_msg)
    except Exception as exc:
        print(f"FAIL: {type(exc).__name__}: {exc}", file=sys.stderr)
        return 3

    print(
        f"OK: wrote {pptx_path} "
        f"({result['slide_count']} slides, {result['text_frames']} text frames, "
        f"fidelity ~{result['fidelity_10']}/10)"
    )
    if not single_file_mode:
        print(f"  report: {report_path}")
        print(f"  json:   {json_path}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
