#!/usr/bin/env python3
"""Deterministic PPTX generator (v2: dark theme + full layout coverage).

Reads a slide-map JSON and produces a .pptx file. No LLM calls -- pure
mechanical slide construction using python-pptx.

Layout coverage (matches `extract_slides.py` output):
    title, content, metric_grid, layer_rows, pipeline, code_block,
    compare_table_2col, compare_table_3col, outcome_grid, split_narrow,
    closing, section/section_divider (legacy passthrough).

Theme: matches `vexjoy-agent-management-deck.html` -- dark navy background,
Aptos body font, Cascadia Code mono, sky-blue accent, semantic colors.
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


# ---------------------------------------------------------------------------
# THEME — single source of truth for v2 visual identity.
# Matches the dark navy HTML deck. Keep colors here; never inline hex codes
# in builders.
# ---------------------------------------------------------------------------
THEME = {
    # Backgrounds
    "bg": RGBColor(0x1A, 0x1A, 0x2E),  # slide background
    "card_bg": RGBColor(0x23, 0x23, 0x40),  # surface / card
    "code_bg": RGBColor(0x16, 0x16, 0x2A),  # code panel
    "border": RGBColor(0x3A, 0x3A, 0x5C),
    # Text
    "fg": RGBColor(0xE8, 0xE8, 0xF0),  # primary text (off-white)
    "fg_sec": RGBColor(0xA0, 0xA0, 0xB8),  # secondary
    "muted": RGBColor(0x6E, 0x6E, 0x8A),  # muted / eyebrow
    # Accents
    "accent": RGBColor(0x64, 0xB5, 0xF6),  # sky blue (brand)
    "success": RGBColor(0x81, 0xC7, 0x84),  # softer green for dark bg
    "danger": RGBColor(0xEF, 0x53, 0x50),  # red
    # Fonts (chained fallbacks aren't supported by python-pptx in a single
    # field; PowerPoint 2023+ ships Aptos and Cascadia Code by default. If
    # absent, PowerPoint falls back to its own substitution table — which
    # gives Segoe UI / Calibri on Windows and Helvetica Neue on macOS.)
    "font_body": "Aptos",
    "font_mono": "Cascadia Code",
}


# ---------------------------------------------------------------------------
# Legacy palette support (kept so old callers don't break, but builders all
# read THEME).
# ---------------------------------------------------------------------------
PALETTES = {
    "vexjoy-dark": {
        "primary": "#1A1A2E",
        "secondary": "#23233D",
        "accent": "#64B5F6",
        "background": "#1A1A2E",
        "text": "#E8E8F0",
        "muted": "#6E6E8A",
    },
}


def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def get_palette(name: str) -> dict:
    raw = PALETTES.get(name.lower(), PALETTES["vexjoy-dark"])
    return {role: hex_to_rgb(color) for role, color in raw.items()}


# ---------------------------------------------------------------------------
# Low-level primitives (theme-aware)
# ---------------------------------------------------------------------------


def fill_bg(slide, color: RGBColor | None = None) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color or THEME["bg"]


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor | None = None,
    align=PP_ALIGN.LEFT,
    font: str | None = None,
    word_wrap: bool = True,
):
    """Add a single-run text box."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font or THEME["font_body"]
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color if color is not None else THEME["fg"]
    return tb


def add_rect(
    slide,
    left,
    top,
    width,
    height,
    *,
    fill_color: RGBColor | None = None,
    line_color: RGBColor | None = None,
    line_width_pt: float = 0.75,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or THEME["card_bg"]
    shape.line.color.rgb = line_color or THEME["border"]
    shape.line.width = Pt(line_width_pt)
    shape.shadow.inherit = False
    shape.text_frame.text = ""
    return shape


def add_eyebrow_and_title(slide, eyebrow: str, title: str, *, title_size: int = 32):
    """Standard top-of-slide pattern: eyebrow label + bold title."""
    if eyebrow:
        add_text(
            slide,
            Inches(0.6),
            Inches(0.5),
            Inches(12),
            Inches(0.4),
            eyebrow.upper(),
            size=12,
            bold=True,
            color=THEME["accent"],
        )
    if title:
        add_text(
            slide,
            Inches(0.6),
            Inches(1.0),
            Inches(12),
            Inches(1.0),
            title,
            size=title_size,
            bold=True,
            color=THEME["fg"],
        )


# ---------------------------------------------------------------------------
# Slide builders. Each takes (prs, slide_data) and returns the slide.
# ---------------------------------------------------------------------------


def _new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill_bg(slide)
    return slide


def build_title(prs, slide_data, _palette=None):
    s = _new_slide(prs)
    eyebrow = slide_data.get("eyebrow", "")
    title = slide_data.get("title", "Untitled")
    subtitle = slide_data.get("subtitle", "")

    if eyebrow:
        add_text(
            s,
            Inches(0.5),
            Inches(0.5),
            Inches(12.3),
            Inches(0.5),
            eyebrow.upper(),
            size=14,
            bold=True,
            color=THEME["accent"],
            align=PP_ALIGN.CENTER,
        )
    add_text(
        s,
        Inches(0.5),
        Inches(2.8),
        Inches(12.3),
        Inches(1.5),
        title,
        size=56,
        bold=True,
        color=THEME["fg"],
        align=PP_ALIGN.CENTER,
    )
    if subtitle:
        add_text(
            s,
            Inches(0.5),
            Inches(4.6),
            Inches(12.3),
            Inches(2.0),
            subtitle,
            size=22,
            color=THEME["fg_sec"],
            align=PP_ALIGN.CENTER,
        )
    return s


def build_content(prs, slide_data, _palette=None):
    """Eyebrow + title + lead paragraph + optional bullets + optional callout."""
    s = _new_slide(prs)
    eyebrow = slide_data.get("eyebrow", "")
    title = slide_data.get("title", "")
    lead = slide_data.get("lead", "")
    callout = slide_data.get("callout", "")
    bullets = slide_data.get("bullets", [])

    add_eyebrow_and_title(s, eyebrow, title, title_size=32)

    y = 2.3
    if lead:
        add_text(
            s,
            Inches(0.6),
            Inches(y),
            Inches(12.0),
            Inches(2.5),
            lead,
            size=18,
            color=THEME["fg_sec"],
        )
        # Approx height per ~110 chars ≈ 0.45in
        y += max(1.0, 0.45 * (len(lead) // 110 + 1))

    if bullets:
        y = max(y, 3.4)
        for b in bullets:
            if isinstance(b, dict):
                prefix = b.get("bold_prefix", "")
                txt = b.get("text", "")
            else:
                prefix, txt = "", str(b)
            tb = s.shapes.add_textbox(Inches(0.9), Inches(y), Inches(11.5), Inches(0.6))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            # Bullet glyph
            r0 = p.add_run()
            r0.text = "• "
            r0.font.name = THEME["font_body"]
            r0.font.size = Pt(18)
            r0.font.color.rgb = THEME["accent"]
            if prefix:
                r1 = p.add_run()
                r1.text = prefix + " "
                r1.font.name = THEME["font_body"]
                r1.font.size = Pt(18)
                r1.font.bold = True
                r1.font.color.rgb = THEME["fg"]
            r2 = p.add_run()
            r2.text = txt
            r2.font.name = THEME["font_body"]
            r2.font.size = Pt(18)
            r2.font.color.rgb = THEME["fg_sec"]
            y += 0.6

    if callout:
        cy = max(y + 0.2, 5.6)
        add_rect(
            s,
            Inches(0.6),
            Inches(cy),
            Inches(12.2),
            Inches(1.4),
            fill_color=THEME["card_bg"],
            line_color=THEME["accent"],
        )
        add_text(
            s,
            Inches(0.85),
            Inches(cy + 0.2),
            Inches(11.7),
            Inches(1.0),
            callout,
            size=14,
            italic=True,
            color=THEME["fg"],
        )
    return s


def build_section_divider(prs, slide_data, _palette=None):
    s = _new_slide(prs)
    add_text(
        s,
        Inches(0.6),
        Inches(0.5),
        Inches(12),
        Inches(0.5),
        slide_data.get("eyebrow", "SECTION").upper(),
        size=14,
        bold=True,
        color=THEME["accent"],
    )
    add_text(
        s,
        Inches(1.0),
        Inches(2.8),
        Inches(11.3),
        Inches(2.0),
        slide_data.get("title", ""),
        size=44,
        bold=True,
        color=THEME["fg"],
        align=PP_ALIGN.LEFT,
    )
    if slide_data.get("subtitle"):
        add_text(
            s,
            Inches(1.0),
            Inches(4.8),
            Inches(11.3),
            Inches(1.5),
            slide_data.get("subtitle", ""),
            size=20,
            color=THEME["fg_sec"],
        )
    return s


def build_metric_grid(prs, slide_data, _palette=None):
    s = _new_slide(prs)
    add_eyebrow_and_title(
        s,
        slide_data.get("eyebrow", ""),
        slide_data.get("title", ""),
    )

    metrics = slide_data.get("metrics", [])
    n = max(1, len(metrics))
    card_w = Inches(2.85) if n >= 4 else Inches(3.5)
    card_h = Inches(2.5)
    gap = Inches(0.18)
    total_w = card_w * n + gap * (n - 1)
    left0 = (prs.slide_width - total_w) / 2
    top0 = Inches(2.6)

    for i, m in enumerate(metrics):
        val = m.get("value", "")
        lab = m.get("label", "")
        desc = m.get("desc", "")
        x = left0 + i * (card_w + gap)
        add_rect(s, x, top0, card_w, card_h)
        add_text(
            s,
            x,
            top0 + Inches(0.4),
            card_w,
            Inches(1.0),
            val,
            size=48,
            bold=True,
            color=THEME["accent"],
            align=PP_ALIGN.CENTER,
        )
        add_text(
            s,
            x,
            top0 + Inches(1.4),
            card_w,
            Inches(0.4),
            lab,
            size=14,
            bold=True,
            color=THEME["muted"],
            align=PP_ALIGN.CENTER,
        )
        if desc:
            add_text(
                s,
                x,
                top0 + Inches(1.85),
                card_w,
                Inches(0.5),
                desc,
                size=12,
                color=THEME["fg_sec"],
                align=PP_ALIGN.CENTER,
            )

    callout = slide_data.get("callout", "")
    if callout:
        add_text(
            s,
            Inches(0.6),
            Inches(5.6),
            Inches(12),
            Inches(1.5),
            callout,
            size=14,
            color=THEME["fg_sec"],
            align=PP_ALIGN.CENTER,
        )
    return s


def build_layer_rows(prs, slide_data, _palette=None):
    s = _new_slide(prs)
    add_eyebrow_and_title(
        s,
        slide_data.get("eyebrow", ""),
        slide_data.get("title", ""),
    )
    layers = slide_data.get("layers", [])
    y = Inches(2.4)
    for layer in layers:
        name = layer.get("name", "")
        count = layer.get("count", "")
        desc = layer.get("desc", "")
        add_rect(s, Inches(0.6), y, Inches(12.2), Inches(0.95))
        add_text(s, Inches(0.9), y + Inches(0.25), Inches(2), Inches(0.5), name, size=22, bold=True, color=THEME["fg"])
        add_text(
            s,
            Inches(3.0),
            y + Inches(0.25),
            Inches(1.2),
            Inches(0.5),
            str(count),
            size=22,
            bold=True,
            color=THEME["accent"],
        )
        add_text(s, Inches(4.5), y + Inches(0.3), Inches(8.2), Inches(0.5), desc, size=14, color=THEME["fg_sec"])
        y += Inches(1.05)
    return s


def build_pipeline(prs, slide_data, _palette=None):
    s = _new_slide(prs)
    add_eyebrow_and_title(
        s,
        slide_data.get("eyebrow", ""),
        slide_data.get("title", ""),
    )
    steps = slide_data.get("pipeline_steps", [])
    n = max(1, len(steps))
    sw = Inches(1.85) if n >= 6 else Inches(2.1)
    sh = Inches(1.2)
    sgap = Inches(0.15)
    total = sw * n + sgap * (n - 1)
    sl = (prs.slide_width - total) / 2
    sy = Inches(3.0)
    for i, step in enumerate(steps):
        label = step.get("label", f"{i + 1:02d}")
        name = step.get("name", "")
        x = sl + i * (sw + sgap)
        add_rect(s, x, sy, sw, sh, line_color=THEME["accent"])
        add_text(
            s,
            x,
            sy + Inches(0.15),
            sw,
            Inches(0.4),
            label,
            size=12,
            bold=True,
            color=THEME["muted"],
            align=PP_ALIGN.CENTER,
        )
        add_text(
            s, x, sy + Inches(0.55), sw, Inches(0.5), name, size=18, bold=True, color=THEME["fg"], align=PP_ALIGN.CENTER
        )

    caption = slide_data.get("pipeline_caption", "") or slide_data.get("callout", "")
    if caption:
        add_text(s, Inches(0.6), Inches(5.0), Inches(12), Inches(1.8), caption, size=16, color=THEME["fg_sec"])
    return s


def build_code_block(prs, slide_data, _palette=None):
    s = _new_slide(prs)
    add_eyebrow_and_title(
        s,
        slide_data.get("eyebrow", ""),
        slide_data.get("title", ""),
    )
    code = slide_data.get("code", "")
    add_rect(
        s,
        Inches(0.8),
        Inches(2.3),
        Inches(11.7),
        Inches(4.6),
        fill_color=THEME["code_bg"],
        line_color=THEME["border"],
    )
    tb = s.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.3), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = THEME["font_mono"]
        r.font.size = Pt(13)
        if line.lstrip().startswith("$"):
            r.font.color.rgb = THEME["muted"]
        elif "✓" in line or "Delivered" in line:
            r.font.color.rgb = THEME["success"]
        elif line.lstrip().startswith(">"):
            r.font.color.rgb = THEME["accent"]
        else:
            r.font.color.rgb = THEME["fg"]
    return s


def _row_color(role: str) -> RGBColor:
    role = (role or "").lower()
    if role == "success":
        return THEME["success"]
    if role == "danger":
        return THEME["danger"]
    if role == "label":
        return THEME["fg"]
    return THEME["fg_sec"]


def _cell_text(cell) -> str:
    if isinstance(cell, dict):
        return cell.get("text", "")
    return str(cell)


def _cell_role(cell) -> str:
    if isinstance(cell, dict):
        return cell.get("role", "")
    return ""


def build_compare_table_2col(prs, slide_data, _palette=None):
    s = _new_slide(prs)
    add_eyebrow_and_title(
        s,
        slide_data.get("eyebrow", ""),
        slide_data.get("title", ""),
    )

    intro = slide_data.get("intro", "")
    if intro:
        add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.6), intro, size=16, color=THEME["fg_sec"])

    tbl = slide_data.get("table", {})
    headers = tbl.get("headers", [])
    rows = tbl.get("rows", [])
    ty = Inches(2.7) if intro else Inches(2.2)
    rh = Inches(0.55)

    # Header row
    if len(headers) >= 2:
        add_text(s, Inches(0.7), ty, Inches(5.5), rh, headers[0], size=14, bold=True, color=THEME["muted"])
        add_text(s, Inches(6.5), ty, Inches(6.3), rh, headers[1], size=14, bold=True, color=THEME["muted"])
        ty += Inches(0.65)

    for row in rows:
        if len(row) < 2:
            continue
        l_text = _cell_text(row[0])
        r_text = _cell_text(row[1])
        l_color = _row_color(_cell_role(row[0]))
        r_color = _row_color(_cell_role(row[1]))
        add_text(s, Inches(0.7), ty, Inches(5.5), rh, l_text, size=14, color=l_color)
        add_text(s, Inches(6.5), ty, Inches(6.3), rh, r_text, size=14, color=r_color)
        ty += Inches(0.65)
    return s


def build_compare_table_3col(prs, slide_data, _palette=None):
    s = _new_slide(prs)
    add_eyebrow_and_title(
        s,
        slide_data.get("eyebrow", ""),
        slide_data.get("title", ""),
        title_size=26,
    )

    intro = slide_data.get("intro", "")
    if intro:
        add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.6), intro, size=14, color=THEME["fg_sec"])

    tbl = slide_data.get("table", {})
    headers = tbl.get("headers", [])
    rows = tbl.get("rows", [])
    ty = Inches(2.4) if not intro else Inches(2.7)

    if len(headers) >= 3:
        add_text(s, Inches(0.7), ty, Inches(4.0), Inches(0.5), headers[0], size=13, bold=True, color=THEME["muted"])
        add_text(s, Inches(4.9), ty, Inches(4.0), Inches(0.5), headers[1], size=13, bold=True, color=THEME["muted"])
        add_text(s, Inches(9.1), ty, Inches(4.0), Inches(0.5), headers[2], size=13, bold=True, color=THEME["muted"])
        ty += Inches(0.7)

    for row in rows:
        if len(row) < 3:
            continue
        add_text(
            s,
            Inches(0.7),
            ty,
            Inches(4.0),
            Inches(0.5),
            _cell_text(row[0]),
            size=13,
            color=_row_color(_cell_role(row[0]) or "label"),
        )
        add_text(
            s,
            Inches(4.9),
            ty,
            Inches(4.0),
            Inches(0.5),
            _cell_text(row[1]),
            size=13,
            color=_row_color(_cell_role(row[1]) or "danger"),
        )
        add_text(
            s,
            Inches(9.1),
            ty,
            Inches(4.0),
            Inches(0.5),
            _cell_text(row[2]),
            size=13,
            color=_row_color(_cell_role(row[2]) or "success"),
        )
        ty += Inches(0.7)
    return s


def build_outcome_grid(prs, slide_data, _palette=None):
    s = _new_slide(prs)
    add_eyebrow_and_title(
        s,
        slide_data.get("eyebrow", ""),
        slide_data.get("title", ""),
    )
    outcomes = slide_data.get("outcomes", [])
    ow = Inches(4.0)
    oh = Inches(1.85)
    ogap = Inches(0.15)
    ox0 = Inches(0.6)
    oy0 = Inches(2.4)
    for i, o in enumerate(outcomes):
        col = i % 3
        row = i // 3
        x = ox0 + col * (ow + ogap)
        y = oy0 + row * (oh + ogap)
        add_rect(s, x, y, ow, oh)
        add_text(
            s,
            x + Inches(0.2),
            y + Inches(0.2),
            ow - Inches(0.4),
            Inches(0.6),
            o.get("heading", ""),
            size=16,
            bold=True,
            color=THEME["fg"],
        )
        add_text(
            s,
            x + Inches(0.2),
            y + Inches(0.85),
            ow - Inches(0.4),
            Inches(1.0),
            o.get("body", ""),
            size=12,
            color=THEME["fg_sec"],
        )
    return s


def build_split_narrow(prs, slide_data, _palette=None):
    s = _new_slide(prs)
    split = slide_data.get("split", {})
    left = split.get("left", {})
    rows = split.get("rows", [])

    eyebrow = left.get("eyebrow", slide_data.get("eyebrow", ""))
    title = left.get("title", slide_data.get("title", ""))
    lead = left.get("lead", slide_data.get("lead", ""))
    callout = left.get("callout", slide_data.get("callout", ""))

    if eyebrow:
        add_text(
            s,
            Inches(0.6),
            Inches(0.5),
            Inches(8),
            Inches(0.4),
            eyebrow.upper(),
            size=12,
            bold=True,
            color=THEME["accent"],
        )
    if title:
        add_text(s, Inches(0.6), Inches(1.0), Inches(7.5), Inches(1.2), title, size=32, bold=True, color=THEME["fg"])
    if lead:
        add_text(s, Inches(0.6), Inches(2.4), Inches(7.5), Inches(2.5), lead, size=16, color=THEME["fg_sec"])
    if callout:
        add_rect(s, Inches(0.7), Inches(5.0), Inches(7.3), Inches(1.4), line_color=THEME["accent"])
        add_text(s, Inches(0.95), Inches(5.3), Inches(6.85), Inches(0.9), callout, size=14, color=THEME["fg"])

    # Right rail: card with name/trigger rows
    add_rect(s, Inches(8.4), Inches(1.0), Inches(4.3), Inches(5.4))
    cy = Inches(1.3)
    for r in rows:
        name = r.get("name", "")
        trigger = r.get("trigger", "")
        add_text(s, Inches(8.6), cy, Inches(2.5), Inches(0.6), name, size=16, bold=True, color=THEME["fg"])
        add_text(
            s,
            Inches(11.0),
            cy,
            Inches(1.6),
            Inches(0.6),
            trigger,
            size=16,
            color=THEME["accent"],
            font=THEME["font_mono"],
        )
        cy += Inches(0.95)
    return s


def build_closing(prs, slide_data, _palette=None):
    s = _new_slide(prs)
    eyebrow = slide_data.get("eyebrow", "")
    title = slide_data.get("title", "Thank You")
    accent_text = slide_data.get("accent_text", "")
    subtitle = slide_data.get("subtitle", "")

    if eyebrow:
        add_text(
            s,
            Inches(0.6),
            Inches(0.5),
            Inches(12),
            Inches(0.4),
            eyebrow.upper(),
            size=14,
            bold=True,
            color=THEME["accent"],
            align=PP_ALIGN.CENTER,
        )

    # Split title from accent_text if accent_text appears within title.
    main_text = title
    if accent_text and accent_text in title:
        main_text = title.replace(accent_text, "").strip()

    add_text(
        s,
        Inches(1.0),
        Inches(2.6),
        Inches(11.3),
        Inches(1.5),
        main_text,
        size=32,
        bold=True,
        color=THEME["fg"],
        align=PP_ALIGN.CENTER,
    )
    if accent_text:
        add_text(
            s,
            Inches(1.0),
            Inches(3.9),
            Inches(11.3),
            Inches(1.5),
            accent_text,
            size=32,
            bold=True,
            color=THEME["accent"],
            align=PP_ALIGN.CENTER,
        )
    if subtitle:
        add_text(
            s,
            Inches(1.0),
            Inches(5.7),
            Inches(11.3),
            Inches(1.5),
            subtitle,
            size=16,
            color=THEME["fg_sec"],
            align=PP_ALIGN.CENTER,
        )
    return s


# ---------------------------------------------------------------------------
# Layout dispatcher
# ---------------------------------------------------------------------------

LAYOUT_BUILDERS = {
    "title": build_title,
    "section": build_section_divider,
    "section_divider": build_section_divider,
    "content": build_content,
    "bullets": build_content,
    "content_bullets": build_content,
    "metric_grid": build_metric_grid,
    "metrics": build_metric_grid,
    "layer_rows": build_layer_rows,
    "pipeline": build_pipeline,
    "code_block": build_code_block,
    "code": build_code_block,
    "compare_table_2col": build_compare_table_2col,
    "compare_table_3col": build_compare_table_3col,
    "outcome_grid": build_outcome_grid,
    "outcomes": build_outcome_grid,
    "split_narrow": build_split_narrow,
    "closing": build_closing,
}


# Public registry so callers (and run-unified.py's report) can ask
# "is this type natively supported?".
SUPPORTED_LAYOUTS = set(LAYOUT_BUILDERS.keys())


def build_presentation(slide_map: list, design: dict, output_path: str) -> str:
    palette_name = design.get("palette", "vexjoy-dark")
    palette = get_palette(palette_name)

    template_path = design.get("template_path")
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in slide_map:
        raw_type = slide_data.get("type", "content")
        slide_type = str(raw_type).lower().strip().replace("-", "_").replace(" ", "_")
        builder = LAYOUT_BUILDERS.get(slide_type, build_content)
        builder(prs, slide_data, palette)

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return str(output)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------


def main():
    parser = argparse.ArgumentParser(description="Generate a PPTX presentation from a slide map JSON.")
    parser.add_argument("--slide-map", required=True)
    parser.add_argument("--design", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    slide_map_path = Path(args.slide_map)
    design_path = Path(args.design)

    if not slide_map_path.exists():
        print(f"ERROR: Slide map not found: {slide_map_path}", file=sys.stderr)
        sys.exit(2)
    if not design_path.exists():
        print(f"ERROR: Design config not found: {design_path}", file=sys.stderr)
        sys.exit(2)

    try:
        with open(slide_map_path) as f:
            slide_map = json.load(f)
        with open(design_path) as f:
            design = json.load(f)
    except json.JSONDecodeError as e:
        print(f"ERROR: Invalid JSON: {e}", file=sys.stderr)
        sys.exit(2)

    if not isinstance(slide_map, list) or not slide_map:
        print("ERROR: Slide map must be a non-empty JSON array", file=sys.stderr)
        sys.exit(2)

    try:
        result = build_presentation(slide_map, design, args.output)
        size = Path(result).stat().st_size
        print(f"SUCCESS: {len(slide_map)} slides -> {result} ({size:,} bytes)")
    except Exception as e:
        print(f"ERROR: Generation failed: {e}", file=sys.stderr)
        sys.exit(3)


if __name__ == "__main__":
    main()
